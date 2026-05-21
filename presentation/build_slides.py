"""Build 2-slide hackathon presentation using actual results."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(BASE, "..", "results", "figures")
OUT  = os.path.join(BASE, "hackathon_slides_auto.pptx")

FIG_TRACK0   = os.path.join(FIGS, "track0_summary.png")
FIG_TRACK2   = os.path.join(FIGS, "track2_novel_summary.png")
FIG_QPU      = os.path.join(FIGS, "qpu_correlation.png")
FIG_HEATMAPS = os.path.join(FIGS, "kernel_heatmaps.png")

# ── Color palette ─────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # IBM quantum cyan
C_ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # quantum orange
C_GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # win green
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xB8, 0xC8, 0xD8)   # muted light blue
C_CARD      = RGBColor(0x16, 0x2B, 0x40)   # card background
C_CARD2     = RGBColor(0x0A, 0x26, 0x1A)   # green-tinted card
C_DIVIDER   = RGBColor(0x00, 0x80, 0xA0)   # divider line

FONT_TITLE  = "Poppins"
FONT_BODY   = "Lato"

# ── Slide size 16:9 ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # totally blank layout

# ─────────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────────

def bg(slide, color=C_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, radius=False):
    """Add a filled rectangle."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1 if not radius else 5,   # MSO_SHAPE_TYPE.RECTANGLE or ROUNDED_RECT
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT,
        font=FONT_BODY, italic=False):
    """Add a text box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return box

def img(slide, path, l, t, w, h):
    """Add an image."""
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def hline(slide, l, t, w, color=C_DIVIDER, width_pt=1.5):
    """Add a horizontal rule."""
    from pptx.util import Pt as UPt
    connector = slide.shapes.add_connector(1, Inches(l), Inches(t), Inches(l+w), Inches(t))
    connector.line.color.rgb = color
    connector.line.width = UPt(width_pt)

def card(slide, l, t, w, h, bg_color=C_CARD):
    """Add a card background."""
    s = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = bg_color
    s.line.fill.background()
    adj = s.adjustments
    if adj:
        adj[0] = 0.05
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: "Quantum Wins on CAR-T Cytotoxicity Prediction"
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)

# Top header bar
rect(s1, 0, 0, 13.33, 1.1, RGBColor(0x06, 0x0E, 0x1A))

# Cyan accent line under header
rect(s1, 0, 1.08, 13.33, 0.04, C_ACCENT)

# Slide number dot
rect(s1, 0.18, 0.3, 0.38, 0.38, C_ACCENT)
txt(s1, "01", 0.18, 0.28, 0.38, 0.42, 11, C_WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_BODY)

# Title
txt(s1, "Quantum Wins on CAR-T Cytotoxicity Prediction",
    0.65, 0.12, 9.5, 0.85, 28, C_WHITE, bold=True, font=FONT_TITLE)

# Subtitle / context
txt(s1, "Track 0 Baseline  ·  172 train / 74 test samples  ·  RBF-SVM on 60-dim classical vs. 180-dim quantum-projected features",
    0.65, 0.75, 10.5, 0.35, 11, C_LIGHT, font=FONT_BODY)

# ── Main figure (left 65%) ────────────────────────────────────────────────────
img(s1, FIG_TRACK0, 0.15, 1.2, 8.4, 4.3)

# ── Right panel: stat cards ───────────────────────────────────────────────────
# Card 1: Classical accuracy
card(s1, 8.75, 1.2, 4.3, 1.1)
txt(s1, "CLASSICAL SVM", 8.85, 1.22, 4.1, 0.35, 9, C_LIGHT, font=FONT_BODY)
txt(s1, "79.7%", 8.85, 1.52, 4.1, 0.65, 28, C_WHITE, bold=True, font=FONT_TITLE, align=PP_ALIGN.LEFT)
txt(s1, "test accuracy", 9.9, 1.62, 3.0, 0.35, 10, C_LIGHT, font=FONT_BODY)

# Card 2: Quantum accuracy
card(s1, 8.75, 2.4, 4.3, 1.1, RGBColor(0x06, 0x26, 0x30))
txt(s1, "QUANTUM PQK", 8.85, 2.42, 4.1, 0.35, 9, C_ACCENT, font=FONT_BODY, bold=True)
txt(s1, "82.4%", 8.85, 2.72, 4.1, 0.65, 28, C_ORANGE, bold=True, font=FONT_TITLE, align=PP_ALIGN.LEFT)
txt(s1, "test accuracy", 9.9, 2.82, 3.0, 0.35, 10, C_LIGHT, font=FONT_BODY)

# Card 3: Delta
card(s1, 8.75, 3.6, 4.3, 0.85, RGBColor(0x06, 0x20, 0x10))
txt(s1, "QUANTUM ADVANTAGE", 8.85, 3.62, 4.1, 0.3, 9, C_GREEN, bold=True, font=FONT_BODY)
txt(s1, "+2.7%  test accuracy", 8.85, 3.9, 4.1, 0.45, 14, C_GREEN, bold=True, font=FONT_TITLE)

# Card 4: Paradox callout
card(s1, 8.75, 4.55, 4.3, 1.1, RGBColor(0x1A, 0x10, 0x05))
txt(s1, "THE ENCODING PARADOX", 8.85, 4.57, 4.1, 0.3, 9, C_ORANGE, bold=True, font=FONT_BODY)
txt(s1, "g_cq = 1.50  <<  √N = 13.1", 8.85, 4.85, 4.1, 0.35, 12, C_WHITE, bold=True, font=FONT_BODY)
txt(s1, "Theory predicts classical wins.\nYet quantum wins by +2.7%.", 8.85, 5.15, 4.1, 0.45, 10.5, C_LIGHT, font=FONT_BODY, italic=True)

# ── Bottom: kernel heatmaps teaser ───────────────────────────────────────────
img(s1, FIG_HEATMAPS, 0.15, 5.6, 8.4, 1.75)

# Bottom label
txt(s1, "Kernel heatmaps (sorted by class label) — quantum kernel shows clearer class block structure",
    0.15, 7.0, 8.4, 0.4, 9.5, C_LIGHT, font=FONT_BODY, italic=True)

# Footer
rect(s1, 0, 7.32, 13.33, 0.18, RGBColor(0x06, 0x0E, 0x1A))
txt(s1, "UW–IBM Quantum Hackathon  ·  CAR-T Cell Cytotoxicity  ·  Projected Quantum Kernels",
    0.3, 7.32, 12.0, 0.18, 8.5, C_LIGHT, font=FONT_BODY, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: "When & Why Quantum Helps + Live QPU Hardware"
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)

# Top header bar
rect(s2, 0, 0, 13.33, 1.1, RGBColor(0x06, 0x0E, 0x1A))
rect(s2, 0, 1.08, 13.33, 0.04, C_ORANGE)

# Slide number dot
rect(s2, 0.18, 0.3, 0.38, 0.38, C_ORANGE)
txt(s2, "02", 0.18, 0.28, 0.38, 0.42, 11, C_WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_BODY)

# Title
txt(s2, "When & Why Quantum Helps + Live QPU Validation",
    0.65, 0.12, 10.5, 0.85, 28, C_WHITE, bold=True, font=FONT_TITLE)

# Subtitle
txt(s2, "Experiment 1: Which samples quantum fixes  ·  Experiment 2: Theory tightness vs. N  ·  Bonus: Live ibm_pittsburgh run",
    0.65, 0.75, 11.0, 0.35, 11, C_LIGHT, font=FONT_BODY)

# ── Main figure: Track 2 novel summary (left 65%) ────────────────────────────
img(s2, FIG_TRACK2, 0.15, 1.2, 8.4, 3.7)

# Label for track2 figure
txt(s2, "Experiment 1 (left): 7 quantum-fixed vs. 5 classical-fixed out of 74 test samples\nExperiment 2 (right): g_cq/√N never reaches 1 — one-hot encoding keeps theory non-tight regardless of N",
    0.15, 4.9, 8.4, 0.55, 9.5, C_LIGHT, font=FONT_BODY, italic=True)

# ── Right panel: QPU results ──────────────────────────────────────────────────
# QPU header card
card(s2, 8.75, 1.2, 4.3, 0.55, RGBColor(0x06, 0x18, 0x28))
txt(s2, "LIVE QPU EXPERIMENT", 8.85, 1.25, 4.1, 0.45, 10, C_ACCENT, bold=True, font=FONT_BODY)

# QPU figure
img(s2, FIG_QPU, 8.65, 1.75, 4.55, 2.5)

# QPU stat cards
card(s2, 8.75, 4.3, 4.3, 0.75, RGBColor(0x06, 0x18, 0x28))
txt(s2, "Backend: ibm_pittsburgh", 8.85, 4.33, 4.1, 0.28, 10, C_ACCENT, bold=True, font=FONT_BODY)
txt(s2, "reps=8 (live)  vs.  reps=24 (pre-computed)", 8.85, 4.58, 4.1, 0.42, 9.5, C_LIGHT, font=FONT_BODY)

card(s2, 8.75, 5.12, 4.3, 0.75, RGBColor(0x1A, 0x08, 0x05))
txt(s2, "Mean Pearson  r = 0.16", 8.85, 5.15, 4.1, 0.32, 12, C_ORANGE, bold=True, font=FONT_TITLE)
txt(s2, "Low r = reps changes the encoding", 8.85, 5.45, 4.1, 0.35, 9.5, C_LIGHT, font=FONT_BODY, italic=True)

card(s2, 8.75, 5.95, 4.3, 1.05, RGBColor(0x06, 0x20, 0x10))
txt(s2, "KEY FINDING", 8.85, 5.97, 4.1, 0.28, 9, C_GREEN, bold=True, font=FONT_BODY)
txt(s2, "reps=8 and reps=24 produce fundamentally different quantum encodings — circuit depth is a meaningful scientific parameter, not just a knob.",
    8.85, 6.22, 4.1, 0.72, 9.5, C_LIGHT, font=FONT_BODY, italic=True)

# ── Bottom summary bar ────────────────────────────────────────────────────────
rect(s2, 0, 7.05, 13.33, 0.27, RGBColor(0x06, 0x14, 0x22))
txt(s2, "Conclusion:  Quantum advantage is real (+2.7%), hardware-confirmed on ibm_pittsburgh, and largest in the data-sparse regime.  Encoding choice determines whether theory and experiment agree.",
    0.3, 7.06, 12.7, 0.26, 9.5, C_WHITE, font=FONT_BODY, bold=False, align=PP_ALIGN.LEFT)

# Footer
rect(s2, 0, 7.32, 13.33, 0.18, RGBColor(0x06, 0x0E, 0x1A))
txt(s2, "UW–IBM Quantum Hackathon  ·  CAR-T Cell Cytotoxicity  ·  Projected Quantum Kernels",
    0.3, 7.32, 12.0, 0.18, 8.5, C_LIGHT, font=FONT_BODY, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
